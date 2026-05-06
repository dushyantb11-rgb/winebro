"""
Generates docs/deck/winebro-deck.pptx from the same content + brand
palette as winebro-deck.html. Run from repo root:

    python docs/deck/generate_pptx.py

Outputs: docs/deck/winebro-deck.pptx (16:9, 12 slides)
Dependencies: python-pptx (pre-installed on dev box)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ------- WineBro brand palette (matches lib/core/theme/app_colors.dart) -----
PAPRIKA      = RGBColor(0x93, 0x00, 0x3C)
PAPRIKA_DARK = RGBColor(0x7A, 0x00, 0x30)
PAPRIKA_DEEP = RGBColor(0x6B, 0x00, 0x2C)
PAPRIKA_ON_DARK = RGBColor(0xD1, 0x4B, 0x7A)
THUNDER      = RGBColor(0x25, 0x21, 0x22)
THUNDER_LT   = RGBColor(0x3A, 0x34, 0x38)
SALEM        = RGBColor(0x0F, 0x80, 0x44)
SALEM_OD     = RGBColor(0x1F, 0xBA, 0x68)
GOLD         = RGBColor(0xC9, 0xA1, 0x4B)
GOLD_WARM    = RGBColor(0xD9, 0xA2, 0x3A)
CREAM        = RGBColor(0xFA, 0xEF, 0xD8)
INK          = RGBColor(0xFA, 0xEF, 0xD8)

PLAYFAIR = "Playfair Display"   # Falls back to system serif if missing
MONT     = "Montserrat"

WIDE = Inches(13.333)
TALL = Inches(7.5)


def add_bg(slide, color):
    """Full-bleed background rectangle."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, WIDE, TALL)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def add_gradient_bg(slide, c1, c2):
    """Two-stop diagonal gradient via two overlapping rectangles. Avoids
    the python-pptx gradient API which is awkward; this gives a clean
    half-and-half look that reads as a gradient on screen."""
    # Base
    rect1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, WIDE, TALL)
    rect1.fill.solid()
    rect1.fill.fore_color.rgb = c1
    rect1.line.fill.background()
    # Overlay at 50% width-tilt
    rect2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, 0, 0, WIDE, TALL)
    rect2.fill.solid()
    rect2.fill.fore_color.rgb = c2
    rect2.line.fill.background()
    rect2.rotation = 180
    return rect1


def add_text(slide, x, y, w, h, text, *,
             font=MONT, size=18, bold=False, italic=False,
             color=CREAM, align=PP_ALIGN.LEFT, spacing=1.2):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_eyebrow(slide, x, y, text, color=CREAM):
    add_text(slide, x, y, Inches(8), Inches(0.4),
             text.upper(), font=MONT, size=10, bold=True, color=color)


def add_pill(slide, x, y, text, fill_color, text_color=CREAM,
             w=Inches(2.7), h=Inches(0.5)):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill_color
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = MONT
    run.font.bold = True
    run.font.size = Pt(10)
    run.font.color.rgb = text_color


def add_brand_tag(slide, x, y, label, fill, ink, w=Inches(1.65)):
    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.55))
    tag.adjustments[0] = 0.18
    tag.fill.solid()
    tag.fill.fore_color.rgb = fill
    tag.line.fill.background()
    tf = tag.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = PLAYFAIR
    run.font.bold = True
    run.font.size = Pt(11)
    run.font.color.rgb = ink


def add_slide_meta(slide, footer, num, color=CREAM):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(7), Inches(0.3),
             footer.upper(), font=MONT, size=8, bold=True,
             color=color, spacing=1.0)
    add_text(slide, Inches(11.4), Inches(7.05), Inches(1.5), Inches(0.3),
             num, font=MONT, size=8, bold=True, color=color,
             align=PP_ALIGN.RIGHT, spacing=1.0)


def add_accent_bar(slide, x, y, color=GOLD_WARM, w=Inches(0.9)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


# ============================================================
# Build deck
# ============================================================
prs = Presentation()
prs.slide_width = WIDE
prs.slide_height = TALL
blank = prs.slide_layouts[6]


# ---- 1. TITLE ----
s = prs.slides.add_slide(blank)
add_bg(s, PAPRIKA_DARK)
add_eyebrow(s, Inches(0.6), Inches(0.6), "WineBro · 2026")
add_text(s, Inches(0.6), Inches(1.2), Inches(11), Inches(2.6),
         "Your Bro for\nevery bottle.",
         font=PLAYFAIR, size=82, bold=True, color=CREAM, spacing=0.95)
add_text(s, Inches(0.6), Inches(4.0), Inches(8), Inches(1.2),
         "India's first homegrown wine, spirits and beer pairing app — "
         "built for biryani, not Bordeaux.",
         font=PLAYFAIR, size=22, italic=True, color=CREAM, spacing=1.4)
add_pill(s, Inches(0.6), Inches(5.6), "v1.1 LAUNCHING", GOLD_WARM, THUNDER)
add_pill(s, Inches(3.5), Inches(5.6), "INDIAN-FIRST PAIRING ENGINE",
         RGBColor(0x47, 0x10, 0x28), CREAM, w=Inches(3.6))
add_pill(s, Inches(7.3), Inches(5.6), "55 SKUs · 66 DISHES · 158 AROMAS",
         RGBColor(0x47, 0x10, 0x28), CREAM, w=Inches(3.9))
add_slide_meta(s, "WineBro · Brand & Investor Deck", "01 / 12")


# ---- 2. THE PROBLEM ----
s = prs.slides.add_slide(blank)
add_bg(s, CREAM)
add_accent_bar(s, Inches(0.6), Inches(0.55), GOLD_WARM)
add_eyebrow(s, Inches(0.6), Inches(0.7), "The problem", color=PAPRIKA)
add_text(s, Inches(0.6), Inches(1.3), Inches(11), Inches(2.2),
         "Western wine apps don't know our food.",
         font=PLAYFAIR, size=58, bold=True, color=THUNDER, spacing=1.0)
add_text(s, Inches(0.6), Inches(3.8), Inches(11.5), Inches(1.5),
         "Open Vivino. Search \"what wine pairs with biryani.\" "
         "Get blank stare.\nSearch \"Mirchi ka Salan.\" Get nothing.\n"
         "Search \"Macher Jhol.\" Get the same Australian Shiraz they "
         "recommend for everything.",
         font=MONT, size=15, color=THUNDER, spacing=1.55)
add_text(s, Inches(0.6), Inches(5.7), Inches(11.5), Inches(1.0),
         "India's wine market is going to 4× by 2030. "
         "Not one app is built for our palate.",
         font=MONT, size=17, bold=True, color=PAPRIKA, spacing=1.4)
add_slide_meta(s, "The problem", "02 / 12", color=THUNDER)


# ---- 3. MARKET SIZE ----
s = prs.slides.add_slide(blank)
add_bg(s, THUNDER)
add_eyebrow(s, Inches(0.6), Inches(0.6), "The market")
add_text(s, Inches(0.6), Inches(1.1), Inches(11), Inches(1.2),
         "The numbers nobody is serving.",
         font=PLAYFAIR, size=48, bold=True, color=CREAM, spacing=1.05)

stats = [
    ("₹14k Cr", "Indian wine market by 2030 (4× from today)", GOLD_WARM),
    ("7%",      "Annual alcobev growth, India",                PAPRIKA_ON_DARK),
    ("0",       "Apps built for the Indian palate",           SALEM_OD),
]
for idx, (n, lbl, col) in enumerate(stats):
    x = Inches(0.6 + idx * 4.2)
    add_text(s, x, Inches(2.7), Inches(4.0), Inches(2.0), n,
             font=PLAYFAIR, size=92, bold=True, color=col, spacing=0.95)
    add_text(s, x, Inches(5.0), Inches(4.0), Inches(0.9), lbl.upper(),
             font=MONT, size=10, bold=True, color=CREAM, spacing=1.4)

add_text(s, Inches(0.6), Inches(6.1), Inches(11.5), Inches(0.8),
         "Indian consumers under-buy wine because they don't know what to "
         "drink with their food. WineBro removes that friction.",
         font=PLAYFAIR, size=18, italic=True, color=CREAM, spacing=1.4)
add_slide_meta(s, "The market", "03 / 12")


# ---- 4. THE PRODUCT ----
s = prs.slides.add_slide(blank)
add_bg(s, PAPRIKA)
add_eyebrow(s, Inches(0.6), Inches(0.6), "The product")
add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(1.2),
         "A pocket sommelier — built for India.",
         font=PLAYFAIR, size=44, bold=True, color=CREAM, spacing=1.05)

features = [
    ("📷", "Scan",       "Snap any bottle. Get tasting notes, pairings, where to buy."),
    ("🍛", "Pair",       "Search any dish. Bro ranks 55 bottles by match %."),
    ("📓", "Journal",    "Quick-log in 15 seconds. Or full-sommelier with photos + voice."),
    ("🤝", "Bro Circle", "What your friends are pouring. Real signals, not algorithms."),
]
for idx, (icon, head, body) in enumerate(features):
    x = Inches(0.6 + idx * 3.05)
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.9), Inches(2.85), Inches(2.6))
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x47, 0x10, 0x28)
    card.line.color.rgb = CREAM
    card.line.width = Pt(0.5)
    card.line.transparency = 0.6
    add_text(s, x + Inches(0.25), Inches(3.05), Inches(1), Inches(0.6),
             icon, font=MONT, size=26, color=CREAM)
    add_text(s, x + Inches(0.25), Inches(3.75), Inches(2.5), Inches(0.5),
             head, font=PLAYFAIR, size=18, bold=True, color=CREAM)
    add_text(s, x + Inches(0.25), Inches(4.25), Inches(2.5), Inches(1.2),
             body, font=MONT, size=10, color=CREAM, spacing=1.45)

add_text(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.7),
         "Free. No subscription wall. Revenue from affiliate, premium tier, brand B2B.",
         font=PLAYFAIR, size=18, italic=True, color=CREAM, spacing=1.4)
add_slide_meta(s, "The product", "04 / 12")


# ---- 5. COVERAGE ----
s = prs.slides.add_slide(blank)
add_bg(s, THUNDER)
add_eyebrow(s, Inches(0.6), Inches(0.6), "Coverage", color=GOLD_WARM)
add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(1.0),
         "Every Indian region is first-class.",
         font=PLAYFAIR, size=42, bold=True, color=CREAM, spacing=1.05)
add_text(s, Inches(0.6), Inches(2.3), Inches(12), Inches(1.3),
         "Bengali. Kashmiri. Hyderabadi. Goan. Rajasthani. Chettinad. Udupi. "
         "North Indian. South Indian. Coastal seafood.",
         font=PLAYFAIR, size=18, italic=True, color=CREAM, spacing=1.4)

brands = [
    ("GROVER ZAMPA", RGBColor(0x5D, 0x0F, 0x1E), CREAM),
    ("SULA",         RGBColor(0x0E, 0x81, 0x75), CREAM),
    ("FRATELLI",     RGBColor(0x6E, 0x12, 0x16), CREAM),
    ("KRSMA",        RGBColor(0x10, 0x10, 0x10), GOLD_WARM),
    ("CHAROSA",      RGBColor(0x4D, 0x5A, 0x2D), CREAM),
    ("BIG BANYAN",   RGBColor(0x2E, 0x5A, 0x2E), CREAM),
    ("AMRUT",        RGBColor(0x8B, 0x5A, 0x2B), CREAM),
    ("PAUL JOHN",    RGBColor(0x13, 0x24, 0x3B), CREAM),
    ("RAMPUR",       RGBColor(0x2A, 0x1A, 0x0F), GOLD),
    ("OLD MONK",     RGBColor(0x3E, 0x24, 0x15), GOLD_WARM),
    ("BIRA 91",      RGBColor(0xE8, 0x5A, 0x1A), CREAM),
    ("KINGFISHER",   RGBColor(0xC8, 0x10, 0x2E), CREAM),
]
gx = 0.6
gy = 4.0
for idx, (lbl, fill, ink) in enumerate(brands):
    col = idx % 6
    row = idx // 6
    add_brand_tag(s, Inches(gx + col * 2.0), Inches(gy + row * 0.75),
                  lbl, fill, ink, w=Inches(1.85))

add_text(s, Inches(0.6), Inches(6.1), Inches(12.5), Inches(0.8),
         "12 Indian brands. 55 SKUs. ₹450 Old Monk to ₹15,000 Penfolds. "
         "Every price tier. Every regional drinker seen.",
         font=MONT, size=12, bold=True, color=CREAM, spacing=1.5)
add_slide_meta(s, "Coverage", "05 / 12")


# ---- 6. AROMA WHEEL / CULTURAL MOAT ----
s = prs.slides.add_slide(blank)
add_bg(s, RGBColor(0x2A, 0x1A, 0x0F))
add_eyebrow(s, Inches(0.6), Inches(0.6), "Cultural moat", color=GOLD_WARM)
add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(1.2),
         "Aam. Jamun. Imli. Tandoor.",
         font=PLAYFAIR, size=46, bold=True, color=CREAM, spacing=1.05)
add_text(s, Inches(0.6), Inches(2.3), Inches(12), Inches(1.4),
         "158 aromas. 26 are India-only — and that's the moat. Western "
         "apps will never match this without a 24-month catch-up.",
         font=PLAYFAIR, size=17, italic=True, color=CREAM, spacing=1.4)

cols = [
    ("Fruity",
     "Aam (Mango)\nJamun (Indian blackberry)\nMunakka\nNimbu"),
    ("Spice",
     "Elaichi (Cardamom)\nHing (Asafoetida)\nSaffron (Kesar)\nMace (Javitri)"),
    ("Smoke & Earth",
     "Tandoor smoke\nAgarbatti (Incense)\nKokum\nImli (Tamarind)"),
]
for idx, (head, body) in enumerate(cols):
    x = Inches(0.6 + idx * 4.2)
    add_text(s, x, Inches(4.3), Inches(4.0), Inches(0.6),
             head, font=PLAYFAIR, size=20, bold=True, color=GOLD_WARM)
    add_text(s, x, Inches(5.0), Inches(4.0), Inches(2.0),
             body, font=MONT, size=13, color=CREAM, spacing=1.55)
add_slide_meta(s, "Cultural moat", "06 / 12")


# ---- 7. DATA ASSETS ----
s = prs.slides.add_slide(blank)
add_bg(s, THUNDER)
add_eyebrow(s, Inches(0.6), Inches(0.6), "Defensibility")
add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(1.0),
         "10 proprietary data assets — collecting from day 1.",
         font=PLAYFAIR, size=34, bold=True, color=CREAM, spacing=1.05)

assets = [
    ("D1", "Pairing matrix",                  "Indian-dish ↔ bottle pairings logged",          "50,000"),
    ("D3", "Indian-context aroma vocabulary", "Vernacular aroma terms users actually use",     "200"),
    ("D5", "Cross-category preference graph", "Wine drinkers who also drink whisky/beer/cocktails", "5,000"),
    ("D6", "Pairing success rate",            "\"Did Bro get it right?\" Yes/Maybe/No feedback", "5,000"),
    ("D7", "Restock behaviour",               "Repeat purchase intent signal",                 "1,000"),
    ("D10", "Occasion-context patterns",      "Home vs restaurant vs bar consumption",         "3,000"),
]

# Header row
hy = 2.5
add_text(s, Inches(0.6), Inches(hy), Inches(2.5), Inches(0.4),
         "ASSET", font=MONT, size=9, bold=True, color=GOLD_WARM, spacing=1.0)
add_text(s, Inches(3.2), Inches(hy), Inches(5.5), Inches(0.4),
         "WHAT WE CAPTURE", font=MONT, size=9, bold=True, color=GOLD_WARM, spacing=1.0)
add_text(s, Inches(10.5), Inches(hy), Inches(2.5), Inches(0.4),
         "YEAR-1 TARGET", font=MONT, size=9, bold=True, color=GOLD_WARM,
         align=PP_ALIGN.RIGHT, spacing=1.0)

for idx, (code, name, what, target) in enumerate(assets):
    y = Inches(3.0 + idx * 0.55)
    add_text(s, Inches(0.6), y, Inches(2.5), Inches(0.5),
             f"{code}  {name}", font=MONT, size=11, bold=True, color=CREAM)
    add_text(s, Inches(3.2), y, Inches(7.0), Inches(0.5),
             what, font=MONT, size=11, color=CREAM)
    add_text(s, Inches(10.5), y, Inches(2.5), Inches(0.5),
             target, font=PLAYFAIR, size=20, bold=True, color=GOLD_WARM,
             align=PP_ALIGN.RIGHT)

add_text(s, Inches(0.6), Inches(6.55), Inches(12.5), Inches(0.7),
         "Every BroCard saved adds a row. Every pairing rated tunes the "
         "engine. Every aroma calibration extends the moat.",
         font=PLAYFAIR, size=14, italic=True, color=CREAM, spacing=1.4)
add_slide_meta(s, "Defensibility", "07 / 12")


# ---- 8. TECH ----
s = prs.slides.add_slide(blank)
add_bg(s, THUNDER)
add_eyebrow(s, Inches(0.6), Inches(0.6), "Built right", color=SALEM_OD)
add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(1.0),
         "Production-shipping. Not a prototype.",
         font=PLAYFAIR, size=42, bold=True, color=CREAM, spacing=1.05)

left_items = [
    "36 PRs merged in 8 production sprints, all per-PR-gated",
    "99/99 unit tests passing — pairing engine, gamification, wrap-up",
    "11 Cloud Functions live (push, recommendations, community signals)",
    "Pairing engine v1.1 — Bayesian-shrunk community feedback consumption",
]
right_items = [
    "Flutter (iOS + Android, single codebase)",
    "Firebase asia-south1 — Firestore, FCM, Storage, Cloud Functions",
    "SHA-256 phone-hash friend discovery — privacy-first",
    "21+ age gate, owner-only journal reads, AA-compliant theme",
]
for idx, item in enumerate(left_items):
    add_text(s, Inches(0.6), Inches(2.5 + idx * 0.7),
             Inches(6.0), Inches(0.7),
             "•  " + item,
             font=MONT, size=13, color=CREAM, spacing=1.45)
for idx, item in enumerate(right_items):
    add_text(s, Inches(7.0), Inches(2.5 + idx * 0.7),
             Inches(6.0), Inches(0.7),
             "•  " + item,
             font=MONT, size=13, color=CREAM, spacing=1.45)

add_text(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.7),
         "Release APK 29-35 MB. 99-test unit suite. "
         "Zero post-merge regressions across 8 sprints.",
         font=PLAYFAIR, size=15, italic=True, color=CREAM, spacing=1.4)
add_slide_meta(s, "Built right", "08 / 12")


# ---- 9. REVENUE ----
s = prs.slides.add_slide(blank)
add_bg(s, RGBColor(0x05, 0x2B, 0x18))
add_eyebrow(s, Inches(0.6), Inches(0.6), "Monetisation")
add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(1.0),
         "Three horizons. One product.",
         font=PLAYFAIR, size=44, bold=True, color=CREAM, spacing=1.05)

horizons = [
    ("M1 — M6",  "Affiliate commerce",
     "Every \"Buy\" tap routes through Flipkart / BigBasket / Amazon partner links. ₹50–150 per conversion. Live now."),
    ("M6 — M12", "Premium tier",
     "Razorpay subscription. AI sommelier chat, advanced engine, no banners. ₹199/mo or ₹1,499/yr."),
    ("M12+",     "B2B brand dashboard",
     "Sula, Grover, Amrut, Bira pay for cohort intel + featured surfacing. ₹50K–2L/quarter per partner."),
]
for idx, (when, head, body) in enumerate(horizons):
    x = Inches(0.6 + idx * 4.2)
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), Inches(4.0), Inches(3.6))
    card.adjustments[0] = 0.05
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x0a, 0x18, 0x10)
    card.line.color.rgb = SALEM_OD
    card.line.width = Pt(0.5)
    card.line.transparency = 0.5
    add_pill(s, x + Inches(0.25), Inches(2.7), when, SALEM_OD, THUNDER,
             w=Inches(1.6))
    add_text(s, x + Inches(0.25), Inches(3.5), Inches(3.5), Inches(0.6),
             head, font=PLAYFAIR, size=20, bold=True, color=CREAM)
    add_text(s, x + Inches(0.25), Inches(4.2), Inches(3.5), Inches(2.0),
             body, font=MONT, size=11, color=CREAM, spacing=1.55)

add_text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.7),
         "Free to use. Three revenue lines unlocked at different scale tiers. "
         "No paywall friction at acquisition.",
         font=PLAYFAIR, size=14, italic=True, color=CREAM, spacing=1.4)
add_slide_meta(s, "Monetisation", "09 / 12")


# ---- 10. TWO ASKS ----
s = prs.slides.add_slide(blank)
# Split background
left_bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, WIDE / 2, TALL)
left_bg.fill.solid(); left_bg.fill.fore_color.rgb = PAPRIKA
left_bg.line.fill.background()
right_bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, WIDE / 2, 0, WIDE / 2, TALL)
right_bg.fill.solid(); right_bg.fill.fore_color.rgb = THUNDER
right_bg.line.fill.background()

# Left side — Brand BD
add_eyebrow(s, Inches(0.6), Inches(0.6), "Brand partner ask")
add_text(s, Inches(0.6), Inches(1.1), Inches(5.8), Inches(2.0),
         "No money.\nJust 30 minutes.",
         font=PLAYFAIR, size=42, bold=True, color=CREAM, spacing=1.0)
add_text(s, Inches(0.6), Inches(3.5), Inches(5.8), Inches(1.5),
         "Soft cross-promotion on your channels. We give you featured "
         "surfacing, push broadcasts on new releases, and a quarterly "
         "anonymized cohort report.",
         font=MONT, size=11, color=CREAM, spacing=1.5)
asks_left = [
    "One Instagram post or newsletter mention per quarter",
    "Optional 30-min tasting-team call to refine pairing bias",
    "Optional co-branded D2C landing page",
]
for idx, item in enumerate(asks_left):
    add_text(s, Inches(0.6), Inches(5.2 + idx * 0.45),
             Inches(5.8), Inches(0.5),
             "•  " + item, font=MONT, size=10, color=CREAM, spacing=1.4)

# Right side — Investor
add_eyebrow(s, Inches(7.0), Inches(0.6), "Investor ask")
add_text(s, Inches(7.0), Inches(1.1), Inches(5.8), Inches(2.0),
         "₹50L.\n18 months.",
         font=PLAYFAIR, size=42, bold=True, color=CREAM, spacing=1.0)
add_text(s, Inches(7.0), Inches(3.5), Inches(5.8), Inches(1.5),
         "Pre-seed runway to 50K MAU + first paid tier launch. "
         "Capital deployed against:",
         font=MONT, size=11, color=CREAM, spacing=1.5)
asks_right = [
    "Photography pipeline (₹5L) — 40 bottle + 60 dish hero photos",
    "BD lead (₹15L/yr) — brand partnerships, influencer seeding",
    "Performance marketing (₹20L) — Meta + Google, Tier-1 metro",
    "Premium tier engineering + AI sommelier (founder-driven)",
]
for idx, item in enumerate(asks_right):
    add_text(s, Inches(7.0), Inches(5.2 + idx * 0.45),
             Inches(5.8), Inches(0.5),
             "•  " + item, font=MONT, size=10, color=CREAM, spacing=1.4)

add_slide_meta(s, "Two asks · One product", "10 / 12")


# ---- 11. WHY NOW ----
s = prs.slides.add_slide(blank)
add_bg(s, CREAM)
add_accent_bar(s, Inches(0.6), Inches(0.55), GOLD_WARM)
add_eyebrow(s, Inches(0.6), Inches(0.7), "Why now", color=PAPRIKA)
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(2.5),
         "India is finally drinking wine.\n"
         "We are the answer to \"but what?\"",
         font=PLAYFAIR, size=46, bold=True, color=THUNDER, spacing=1.0)

# Quote bar
qbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.6),
                          Inches(0.05), Inches(2.0))
qbar.fill.solid(); qbar.fill.fore_color.rgb = GOLD_WARM
qbar.line.fill.background()
add_text(s, Inches(0.85), Inches(4.6), Inches(11.5), Inches(2.4),
         "Indian wine consumption tripled in the last decade. Tier-1 metros "
         "lead, but Tier-2 cities are growing faster. The bottleneck for "
         "every new drinker is the same: they don't know what to buy or "
         "what to drink it with. WineBro is the answer.",
         font=PLAYFAIR, size=18, italic=True, color=THUNDER, spacing=1.45)

add_slide_meta(s, "Why now", "11 / 12", color=THUNDER)


# ---- 12. CLOSING ----
s = prs.slides.add_slide(blank)
add_bg(s, PAPRIKA_DARK)
add_eyebrow(s, Inches(0.6), Inches(0.6), "Let's pour something")
add_text(s, Inches(0.6), Inches(1.2), Inches(11), Inches(2.6),
         "Your Bro for\nevery bottle.",
         font=PLAYFAIR, size=82, bold=True, color=CREAM, spacing=0.95)
add_text(s, Inches(0.6), Inches(4.0), Inches(8), Inches(1.0),
         "India's first homegrown wine pairing app — built for biryani, not Bordeaux.",
         font=PLAYFAIR, size=18, italic=True, color=CREAM, spacing=1.4)

contacts = [
    ("App",            "winebro.web.app"),
    ("BD + investor",  "[your email]"),
    ("Phone",          "[your number]"),
]
for idx, (k, v) in enumerate(contacts):
    add_text(s, Inches(0.6), Inches(5.2 + idx * 0.45),
             Inches(8), Inches(0.4),
             k + " ·  " + v, font=MONT, size=12, color=CREAM, spacing=1.4)

add_slide_meta(s, "Drink responsibly · 21+ only", "12 / 12")


# ============================================================
output = Path(__file__).parent / "winebro-deck.pptx"
prs.save(str(output))
print(f"Wrote {output}")
print(f"Slides: {len(prs.slides)}  ·  Aspect: 16:9  ·  Brand-aligned to "
      f"WineBro paprika/thunder/gold palette.")
